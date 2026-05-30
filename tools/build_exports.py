#!/usr/bin/env python3
import os, re, glob
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
EXPORTS = ROOT / "exports"
EXPORTS.mkdir(exist_ok=True)

DOMAIN_TITLES = {
    "01-monetization-digital-empires": "01 - Monetization & Digital Empires",
    "02-media-production": "02 - Media & Production",
    "03-computer-science-architecture": "03 - Computer Science & Architecture",
    "04-security": "04 - Offensive/Defensive Security",
    "05-human-physiology-optimization": "05 - Human Physiology & Optimization",
    "06-deep-sciences-biology": "06 - Deep Sciences & Biology",
    "07-engineering-mechanics": "07 - Engineering & Mechanics",
    "08-combat-kinesiology": "08 - Combat & Kinesiology",
    "09-humanities-master-skills": "09 - Humanities & Master Skills",
    "10-curated-wisdom": "10 - Curated Wisdom",
}

def domain_dirs():
    return sorted(d for d in os.listdir(ROOT)
                  if (ROOT / d).is_dir() and re.match(r"^\d\d-", d))

def md_files(domain):
    return sorted(glob.glob(str(ROOT / domain / "*.md")))

def strip_frontmatter(text):
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            return text[end + 4:].lstrip("\n")
    return text

def parse_blocks(body):
    lines = body.splitlines(); i = 0; in_code = False; code_buf = []
    while i < len(lines):
        line = lines[i]
        if line.strip().startswith("```"):
            if in_code:
                yield ("code", "\n".join(code_buf)); code_buf = []; in_code = False
            else:
                in_code = True
            i += 1; continue
        if in_code:
            code_buf.append(line); i += 1; continue
        if not line.strip():
            i += 1; continue
        if line.startswith("# "):
            yield ("h1", line[2:].strip())
        elif line.startswith("## "):
            yield ("h2", line[3:].strip())
        elif line.startswith("### "):
            yield ("h3", line[4:].strip())
        elif line.startswith("#### "):
            yield ("h3", line[5:].strip())
        elif line.lstrip().startswith(("- ", "* ")):
            yield ("bullet", line.lstrip()[2:].strip())
        elif re.match(r"^\d+\.\s", line.strip()):
            yield ("bullet", re.sub(r"^\d+\.\s", "", line.strip()))
        elif line.startswith(">"):
            yield ("quote", line.lstrip("> ").strip())
        elif line.strip().startswith("|"):
            yield ("table", line.strip())
        else:
            yield ("para", line.strip())
        i += 1

def clean_inline(t):
    t = re.sub(r"`([^`]*)`", r"\1", t)
    t = re.sub(r"\*\*([^*]*)\*\*", r"\1", t)
    t = re.sub(r"\*([^*]*)\*", r"\1", t)
    t = re.sub(r"\[([^\]]*)\]\(([^)]*)\)", r"\1 (\2)", t)
    return t

def build_docx():
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = Document()
    t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("KNOWLEDGE VAULT"); r.bold = True; r.font.size = Pt(34)
    s = doc.add_paragraph(); s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rs = s.add_run("A Multi-Disciplinary Technical Reference - 62 files, 10 domains")
    rs.italic = True; rs.font.size = Pt(13)
    doc.add_page_break()
    for domain in domain_dirs():
        doc.add_heading(DOMAIN_TITLES.get(domain, domain), level=0)
        for path in md_files(domain):
            body = strip_frontmatter(Path(path).read_text(encoding="utf-8"))
            for kind, content in parse_blocks(body):
                c = clean_inline(content)
                if kind == "h1": doc.add_heading(c, level=1)
                elif kind == "h2": doc.add_heading(c, level=2)
                elif kind == "h3": doc.add_heading(c, level=3)
                elif kind == "bullet": doc.add_paragraph(c, style="List Bullet")
                elif kind == "code":
                    p = doc.add_paragraph(); run = p.add_run(content)
                    run.font.name = "Courier New"; run.font.size = Pt(8)
                elif kind == "quote":
                    p = doc.add_paragraph(c); p.style = "Quote"
                else: doc.add_paragraph(c)
            doc.add_page_break()
    out = EXPORTS / "Knowledge-Vault-Complete.docx"; doc.save(str(out)); return out

def build_pdf():
    from fpdf import FPDF
    class PDF(FPDF):
        def footer(self):
            self.set_y(-12); self.set_font("Helvetica", "I", 7)
            self.set_text_color(150,150,150)
            self.cell(0, 8, f"Knowledge Vault - page {self.page_no()}", align="C")
    pdf = PDF(format="A4"); pdf.set_margins(15, 15, 15); pdf.set_auto_page_break(auto=True, margin=16)
    def w(s): return s.encode("latin-1","replace").decode("latin-1")
    pdf.add_page(); pdf.ln(60)
    pdf.set_font("Helvetica","B",30); pdf.multi_cell(pdf.epw,14,w("KNOWLEDGE VAULT"),align="C")
    pdf.set_font("Helvetica","I",13)
    pdf.multi_cell(pdf.epw,9,w("A Multi-Disciplinary Technical Reference\n62 files - 10 domains"),align="C")
    for domain in domain_dirs():
        pdf.add_page(); pdf.set_font("Helvetica","B",20); pdf.set_text_color(20,20,90)
        pdf.multi_cell(pdf.epw,11,w(DOMAIN_TITLES.get(domain,domain))); pdf.set_text_color(0,0,0)
        for path in md_files(domain):
            body = strip_frontmatter(Path(path).read_text(encoding="utf-8")); pdf.ln(2)
            for kind, content in parse_blocks(body):
                c = w(clean_inline(content))
                if kind == "h1": pdf.ln(2); pdf.set_font("Helvetica","B",16); pdf.multi_cell(pdf.epw,8,c)
                elif kind == "h2":
                    pdf.ln(1); pdf.set_font("Helvetica","B",13); pdf.set_text_color(20,20,90)
                    pdf.multi_cell(pdf.epw,7,c); pdf.set_text_color(0,0,0)
                elif kind == "h3": pdf.set_font("Helvetica","B",11); pdf.multi_cell(pdf.epw,6,c)
                elif kind == "bullet":
                    pdf.set_font("Helvetica","",10); pdf.multi_cell(pdf.epw,5.5,w("  - "+clean_inline(content)))
                elif kind == "code":
                    pdf.set_font("Courier","",8); pdf.set_fill_color(244,244,244)
                    for ln in content.splitlines(): pdf.multi_cell(pdf.epw,4.2,w(ln),fill=True)
                    pdf.set_fill_color(255,255,255)
                elif kind == "quote":
                    pdf.set_font("Helvetica","I",10); pdf.multi_cell(pdf.epw,5.5,w("  "+clean_inline(content)))
                elif kind == "table":
                    pdf.set_font("Courier","",8); pdf.multi_cell(pdf.epw,4.5,c)
                else:
                    pdf.set_font("Helvetica","",10); pdf.multi_cell(pdf.epw,5.5,c)
    out = EXPORTS / "Knowledge-Vault-Complete.pdf"; pdf.output(str(out)); return out

def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]; title_layout = prs.slide_layouts[0]
    def first_h1_and_intro(body):
        h1, intro = None, None
        for kind, content in parse_blocks(body):
            if kind == "h1" and h1 is None: h1 = clean_inline(content)
            elif kind == "para" and intro is None and h1 is not None:
                intro = clean_inline(content); break
        return h1, intro
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = "Knowledge Vault"
    s.placeholders[1].text = "A Multi-Disciplinary Technical Reference - 62 files across 10 domains"
    for domain in domain_dirs():
        ds = prs.slides.add_slide(blank)
        box = ds.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(12), Inches(2))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = DOMAIN_TITLES.get(domain, domain)
        p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = RGBColor(0x14,0x14,0x5A)
        files = md_files(domain)
        sub = tf.add_paragraph(); sub.text = f"{len(files)} topics"
        sub.font.size = Pt(20); sub.font.color.rgb = RGBColor(0x66,0x66,0x66)
        for path in files:
            body = strip_frontmatter(Path(path).read_text(encoding="utf-8"))
            h1, intro = first_h1_and_intro(body)
            cs = prs.slides.add_slide(blank)
            tb = cs.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(6.5))
            tf = tb.text_frame; tf.word_wrap = True
            tp = tf.paragraphs[0]; tp.text = h1 or Path(path).stem
            tp.font.size = Pt(28); tp.font.bold = True; tp.font.color.rgb = RGBColor(0x14,0x14,0x5A)
            if intro:
                ip = tf.add_paragraph(); ip.text = intro[:600]; ip.font.size = Pt(15)
            secs = [clean_inline(c) for k, c in parse_blocks(body) if k == "h2"][:6]
            if secs:
                hp = tf.add_paragraph(); hp.text = "Sections:"; hp.font.size = Pt(13); hp.font.bold = True
                for sname in secs:
                    bp = tf.add_paragraph(); bp.text = "- " + sname; bp.font.size = Pt(12)
    out = EXPORTS / "Knowledge-Vault-Overview.pptx"; prs.save(str(out)); return out

if __name__ == "__main__":
    print("Building DOCX..."); print("  ->", build_docx())
    print("Building PDF...");  print("  ->", build_pdf())
    print("Building PPTX..."); print("  ->", build_pptx())
    print("Done. Exports in:", EXPORTS)
